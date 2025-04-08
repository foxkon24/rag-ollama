# onedrive_search.py - OneDriveファイル検索機能（Python実装改善版）
import os
import logging
import re
import time
import json
import glob
from datetime import datetime
import io
from pypdf import PdfReader  # pypdfライブラリ
from docx import Document
import openpyxl
from pptx import Presentation
import chardet

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

class OneDriveSearch:
    def __init__(self, base_directory=None, file_types=None, max_results=10):
        """
        OneDrive検索クラスの初期化

        Args:
            base_directory: 検索の基準ディレクトリ（指定がない場合はOneDriveルート）
            file_types: 検索対象のファイル拡張子リスト
            max_results: デフォルトの最大検索結果数
        """
        # OneDriveのルートディレクトリを取得（環境に応じて調整が必要）
        self.onedrive_root = os.path.expanduser("~/OneDrive")

        # ユーザー名を取得
        self.username = os.getenv("USERNAME", "owner")

        if not os.path.exists(self.onedrive_root):
            # 標準的なOneDriveパスが見つからない場合は代替パスを試す
            alt_paths = [
                os.path.expanduser("~/OneDrive - Company"),  # 企業アカウント用
                os.path.expanduser("~/OneDrive - Personal"),  # 個人アカウント用
                f"C:\\Users\\{self.username}\\OneDrive",  # 絶対パス
                f"D:\\OneDrive",  # 別ドライブ
                # 共立電機製作所のパターンを追加
                f"C:\\Users\\{self.username}\\OneDrive - 株式会社　共立電機製作所"
            ]

            for path in alt_paths:
                if os.path.exists(path):
                    self.onedrive_root = path
                    break

        logger.info(f"OneDriveルートディレクトリ: {self.onedrive_root}")

        # 検索の基準ディレクトリを設定
        self.base_directory = base_directory if base_directory else self.onedrive_root
        logger.info(f"検索基準ディレクトリ: {self.base_directory}")

        # ファイルタイプの設定
        self.file_types = file_types if file_types else []
        logger.info(f"検索対象ファイルタイプ: {', '.join(self.file_types) if self.file_types else '全ファイル'}")

        # 最大結果数
        self.max_results = max_results
        logger.info(f"デフォルト最大検索結果数: {self.max_results}")

        # 検索結果キャッシュ（パフォーマンス向上のため）
        self.search_cache = {}
        self.cache_expiry = 300  # キャッシュの有効期限（秒）

    def search_files(self, keywords, file_types=None, max_results=None, use_cache=True):
        """
        OneDrive内のファイルをキーワードで検索（Pythonネイティブ実装）

        Args:
            keywords: 検索キーワード（文字列またはリスト）
            file_types: 検索対象の拡張子リスト（例: ['.pdf', '.docx']）
            max_results: 最大結果数
            use_cache: キャッシュを使用するかどうか

        Returns:
            検索結果のリスト [{'path': ファイルパス, 'name': ファイル名, 'modified': 更新日時}]
        """
        # デフォルト値の設定
        if file_types is None:
            file_types = self.file_types

        if max_results is None:
            max_results = self.max_results

        # キャッシュキーの生成
        cache_key = f"{str(keywords)}_{str(file_types)}_{max_results}"

        # キャッシュチェック
        if use_cache and cache_key in self.search_cache:
            cache_entry = self.search_cache[cache_key]
            cache_time = cache_entry['timestamp']
            current_time = time.time()

            # キャッシュが有効期限内なら使用
            if current_time - cache_time < self.cache_expiry:
                logger.info(f"キャッシュから検索結果を返します: {len(cache_entry['results'])}件")
                return cache_entry['results']

        # キーワードを文字列から配列に変換
        if isinstance(keywords, str):
            keywords = keywords.split()

        # 日報関連の特別キーワードを抽出
        date_keywords = []
        search_terms = []

        for k in keywords:
            # 日付形式（YYYY年MM月DD日）を抽出
            if re.search(r'\d{4}年\d{1,2}月\d{1,2}日', k):
                date_match = re.search(r'(\d{4})年(\d{1,2})月(\d{1,2})日', k)
                if date_match:
                    year = date_match.group(1)
                    month = date_match.group(2).zfill(2)  # 1桁の月を2桁に
                    day = date_match.group(3).zfill(2)    # 1桁の日を2桁に
                    date_pattern = f"{year}{month}{day}"
                    date_pattern2 = f"{year}-{month}-{day}"
                    date_pattern3 = f"{year}/{month}/{day}"
                    date_keywords.extend([date_pattern, date_pattern2, date_pattern3])
            else:
                # 日本語検索キーワードは短くして検索精度を上げる
                if len(k) > 2 and re.search(r'[ぁ-んァ-ン一-龥]', k):
                    search_terms.append(k)

        # 少なくとも日付キーワードは追加
        if date_keywords:
            search_terms.extend(date_keywords)

        # 検索キーワードがない場合、元のキーワードの先頭2つを使用
        if not search_terms and keywords:
            search_terms = keywords[:2]

        # 最低1つのキーワードを確保
        if not search_terms and isinstance(keywords, str) and keywords:
            search_terms = [keywords]

        logger.info(f"OneDrive検索を実行: キーワード={search_terms}, ファイルタイプ={file_types}")

        try:
            # Pythonネイティブの実装でファイル検索を行う
            search_path = self.base_directory
            results = []
            
            # 日付パターンがあるか確認
            has_date_pattern = bool(date_keywords)
            
            # ファイル検索のベースパターン
            search_pattern = "**/*"  # 再帰的に全ファイル検索
            
            # 日付フォルダパターンを構築（日付キーワードがある場合）
            date_path_patterns = []
            if has_date_pattern and date_keywords:
                for date_key in date_keywords:
                    if len(date_key) == 8 and date_key.isdigit():  # YYYYMMDD形式
                        year = date_key[:4]
                        month = date_key[4:6]
                        day = date_key[6:8]
                        # 様々な日付フォルダパターンを試す
                        date_path_patterns.extend([
                            f"**/{year}*/**/{month}*/**/{day}*/**",
                            f"**/{year}*/{month}*/{day}*/**",
                            f"**/{year}年度/**",
                            f"**/{year}/**",
                        ])
            
            # 検索パスリストを作成
            search_paths = [search_pattern]
            if date_path_patterns:
                search_paths.extend(date_path_patterns)
            
            # 各ファイルタイプごとに検索
            all_matched_files = []
            
            if not file_types:
                # ファイルタイプ指定がなければ全ファイル検索
                for pattern in search_paths:
                    matched_files = glob.glob(os.path.join(search_path, pattern), recursive=True)
                    all_matched_files.extend(matched_files)
            else:
                # 指定されたファイルタイプごとに検索
                for ext in file_types:
                    for pattern in search_paths:
                        # ファイル拡張子を追加
                        if not pattern.endswith("*"):
                            file_pattern = f"{pattern}*{ext}"
                        else:
                            file_pattern = f"{pattern}{ext}"
                        
                        matched_files = glob.glob(os.path.join(search_path, file_pattern), recursive=True)
                        all_matched_files.extend(matched_files)
            
            # 重複を排除
            all_matched_files = list(set(all_matched_files))
            
            # キーワード検索: ファイル名にキーワードが含まれるものをフィルタリング
            filtered_files = []
            for file_path in all_matched_files:
                if not os.path.isfile(file_path):
                    continue
                    
                file_name = os.path.basename(file_path)
                matches_keywords = False
                
                # 日付キーワードがある場合
                if date_keywords:
                    for date_key in date_keywords:
                        if date_key in file_name or date_key.replace("-", "") in file_name or date_key.replace("/", "") in file_name:
                            matches_keywords = True
                            break
                
                # その他のキーワード
                if not matches_keywords and search_terms:
                    other_terms = [term for term in search_terms if term not in date_keywords]
                    if other_terms:
                        for term in other_terms:
                            if term in file_name:
                                matches_keywords = True
                                break
                
                # 日付のみで検索する場合やキーワードがない場合は全ファイルを対象に
                if (not search_terms) or (has_date_pattern and not other_terms) or matches_keywords:
                    filtered_files.append(file_path)
            
            # ファイルのメタデータを取得
            for file_path in filtered_files[:max_results]:
                try:
                    stat_info = os.stat(file_path)
                    file_name = os.path.basename(file_path)
                    
                    results.append({
                        'path': file_path,
                        'name': file_name,
                        'modified': datetime.fromtimestamp(stat_info.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
                        'size': stat_info.st_size,
                        'extension': os.path.splitext(file_name)[1].lower()
                    })
                except (FileNotFoundError, PermissionError) as e:
                    logger.warning(f"ファイル {file_path} の情報取得中にエラー: {str(e)}")
            
            # 結果のフォーマットと表示
            logger.info(f"検索結果: {len(results)}件")
            for i, result in enumerate(results[:3]):  # 最初の3件のみログ表示
                logger.info(f"結果{i+1}: {result.get('name')} - {result.get('path')}")

            # キャッシュに保存
            self.search_cache[cache_key] = {
                'results': results,
                'timestamp': time.time()
            }

            return results

        except Exception as e:
            logger.error(f"OneDrive検索中にエラーが発生しました: {str(e)}")
            logger.error(f"詳細: {str(e.__class__.__name__)}")
            return []

    def detect_encoding(self, file_path):
        """
        ファイルのエンコーディングを検出する

        Args:
            file_path: ファイルパス

        Returns:
            検出されたエンコーディング
        """
        try:
            with open(file_path, 'rb') as f:
                result = chardet.detect(f.read(10000))  # 最初の10000バイトでエンコーディングを推測
            return result['encoding'] or 'utf-8'
        except Exception as e:
            logger.warning(f"エンコーディング検出エラー: {str(e)}")
            return 'utf-8'

    def read_file_content(self, file_path):
        """
        ファイルの内容を読み込む（Pythonライブラリ使用）

        Args:
            file_path: 読み込むファイルパス

        Returns:
            ファイルの内容（文字列）
        """
        try:
            # ファイルの拡張子を取得
            _, ext = os.path.splitext(file_path.lower())

            # テキストファイル形式かどうか判定
            text_extensions = ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.log', '.py', '.js', '.css']
            
            if ext in text_extensions:
                # テキストファイルの場合
                encoding = self.detect_encoding(file_path)
                with open(file_path, 'r', encoding=encoding, errors='replace') as f:
                    content = f.read()
                logger.info(f"テキストファイルとして読み込みました ({encoding}): {file_path}")
                return content
                
            elif ext == '.pdf':
                return self._extract_pdf_content_python(file_path)
                
            elif ext == '.docx':
                return self._extract_word_content_python(file_path)
                
            elif ext == '.xlsx':
                return self._extract_excel_content_python(file_path)
                
            elif ext == '.pptx':
                return self._extract_powerpoint_content_python(file_path)
                
            else:
                # サポートされていないファイル形式
                logger.warning(f"サポートされていないファイル形式です: {ext}")
                file_info = os.stat(file_path)
                size_kb = file_info.st_size / 1024
                modified = datetime.fromtimestamp(file_info.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
                
                return f"""ファイル: {os.path.basename(file_path)}
サイズ: {size_kb:.2f} KB
更新日時: {modified}
---------------------------------
このファイル形式({ext})は内容の読み込みに対応していません。"""

        except PermissionError:
            logger.error(f"ファイル '{file_path}' へのアクセス権限がありません")
            return f"ファイル '{os.path.basename(file_path)}' へのアクセス権限がありません。システム管理者に確認してください。"
        except FileNotFoundError:
            logger.error(f"ファイル '{file_path}' が見つかりません")
            return f"ファイル '{os.path.basename(file_path)}' が見つかりません。削除または移動された可能性があります。"
        except Exception as e:
            logger.error(f"ファイル読み込み中にエラーが発生しました: {str(e)}")
            return f"ファイル読み込みエラー: {str(e)}"

    def _extract_pdf_content_python(self, file_path):
        """
        PDFファイルからテキストを抽出する（pypdfライブラリ使用）
        """
        try:
            pdf_text = []
            # ファイル情報を追加
            file_info = os.stat(file_path)
            size_kb = file_info.st_size / 1024
            modified = datetime.fromtimestamp(file_info.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            
            pdf_text.append(f"PDF名: {os.path.basename(file_path)}")
            pdf_text.append(f"ファイルサイズ: {size_kb:.2f} KB")
            pdf_text.append(f"最終更新日時: {modified}")
            pdf_text.append("----------------------------------------")
            
            # PDFの内容を抽出
            reader = PdfReader(file_path)
            
            # ドキュメント情報を追加
            pdf_text.append(f"ページ数: {len(reader.pages)}")
            
            # メタデータを追加
            if reader.metadata:
                pdf_text.append("メタデータ:")
                for key, value in reader.metadata.items():
                    if value and not key.startswith('/'):
                        key_name = key[1:] if key.startswith('/') else key
                        pdf_text.append(f"  {key_name}: {value}")
            
            # 各ページからテキストを抽出（最初の3ページまたは全ページ）
            max_pages = min(len(reader.pages), 3)
            for page_num in range(max_pages):
                page = reader.pages[page_num]
                pdf_text.append(f"\n[ページ {page_num+1}]")
                text = page.extract_text()
                # 長いテキストは要約
                if text and len(text) > 1000:
                    text = text[:1000] + "...(続く)"
                pdf_text.append(text if text else "ページにテキストが見つかりません")
            
            if len(reader.pages) > max_pages:
                pdf_text.append(f"\n...(残り {len(reader.pages) - max_pages} ページは省略されました)")
            
            return "\n".join(pdf_text)
            
        except Exception as e:
            logger.error(f"PDF抽出エラー: {str(e)}")
            return f"PDF '{os.path.basename(file_path)}' の読み込みに失敗しました: {str(e)}"

    def _extract_word_content_python(self, file_path):
        """
        Wordファイル(docx)からテキストを抽出する（python-docx使用）
        """
        try:
            word_text = []
            # ファイル情報を追加
            file_info = os.stat(file_path)
            size_kb = file_info.st_size / 1024
            modified = datetime.fromtimestamp(file_info.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            
            word_text.append(f"Word文書名: {os.path.basename(file_path)}")
            word_text.append(f"ファイルサイズ: {size_kb:.2f} KB")
            word_text.append(f"最終更新日時: {modified}")
            word_text.append("----------------------------------------")
            
            # Wordファイルを開く
            doc = Document(file_path)
            
            # 文書のプロパティ
            core_props = doc.core_properties
            if core_props:
                word_text.append("ドキュメント情報:")
                if core_props.title:
                    word_text.append(f"  タイトル: {core_props.title}")
                if core_props.author:
                    word_text.append(f"  作成者: {core_props.author}")
                if core_props.created:
                    word_text.append(f"  作成日時: {core_props.created}")
                if core_props.modified:
                    word_text.append(f"  更新日時: {core_props.modified}")
            
            # 段落を抽出
            word_text.append("\n文書内容:")
            paragraphs_text = []
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    paragraphs_text.append(para.text)
                # 長すぎる場合は途中で切る
                if i > 50 or len("\n".join(paragraphs_text)) > 3000:
                    paragraphs_text.append("...(以下省略)")
                    break
            
            word_text.append("\n".join(paragraphs_text))
            
            return "\n".join(word_text)
            
        except Exception as e:
            logger.error(f"Word抽出エラー: {str(e)}")
            return f"Word文書 '{os.path.basename(file_path)}' の読み込みに失敗しました: {str(e)}"

    def _extract_excel_content_python(self, file_path):
        """
        Excelファイル(xlsx)から内容を抽出する（openpyxl使用）
        """
        try:
            excel_text = []
            # ファイル情報を追加
            file_info = os.stat(file_path)
            size_kb = file_info.st_size / 1024
            modified = datetime.fromtimestamp(file_info.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            
            excel_text.append(f"Excel名: {os.path.basename(file_path)}")
            excel_text.append(f"ファイルサイズ: {size_kb:.2f} KB")
            excel_text.append(f"最終更新日時: {modified}")
            excel_text.append("----------------------------------------")
            
            # Excelファイルを開く
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            # ブック情報
            excel_text.append(f"シート数: {len(wb.sheetnames)}")
            excel_text.append(f"シート一覧: {', '.join(wb.sheetnames)}")
            
            # 各シートの内容を抽出（最初の2シートまで）
            sheet_count = 0
            for sheet_name in wb.sheetnames:
                if sheet_count >= 2:
                    excel_text.append(f"\n...(残り {len(wb.sheetnames) - 2} シートは省略されました)")
                    break
                
                sheet = wb[sheet_name]
                excel_text.append(f"\n[シート: {sheet_name}]")
                
                # シートの次元
                if sheet.max_row > 0 and sheet.max_column > 0:
                    excel_text.append(f"データ範囲: A1:{openpyxl.utils.get_column_letter(sheet.max_column)}{sheet.max_row}")
                
                # データの抽出（最初の10行まで）
                max_rows = min(sheet.max_row, 10)
                max_cols = min(sheet.max_column, 10)
                
                for row in range(1, max_rows + 1):
                    row_data = []
                    for col in range(1, max_cols + 1):
                        cell_value = sheet.cell(row=row, column=col).value
                        if cell_value is not None:
                            row_data.append(str(cell_value))
                        else:
                            row_data.append("")
                    excel_text.append(" | ".join(row_data))
                
                if sheet.max_row > max_rows:
                    excel_text.append(f"...(残り {sheet.max_row - max_rows} 行は省略されました)")
                
                if sheet.max_column > max_cols:
                    excel_text.append(f"...(残り {sheet.max_column - max_cols} 列は省略されました)")
                
                sheet_count += 1
            
            wb.close()
            return "\n".join(excel_text)
            
        except Exception as e:
            logger.error(f"Excel抽出エラー: {str(e)}")
            return f"Excelファイル '{os.path.basename(file_path)}' の読み込みに失敗しました: {str(e)}"

    def _extract_powerpoint_content_python(self, file_path):
        """
        PowerPointファイル(pptx)から内容を抽出する（python-pptx使用）
        """
        try:
            ppt_text = []
            # ファイル情報を追加
            file_info = os.stat(file_path)
            size_kb = file_info.st_size / 1024
            modified = datetime.fromtimestamp(file_info.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            
            ppt_text.append(f"PowerPoint名: {os.path.basename(file_path)}")
            ppt_text.append(f"ファイルサイズ: {size_kb:.2f} KB")
            ppt_text.append(f"最終更新日時: {modified}")
            ppt_text.append("----------------------------------------")
            
            # PowerPointファイルを開く
            prs = Presentation(file_path)
            
            # プレゼンテーション情報
            ppt_text.append(f"スライド数: {len(prs.slides)}")
            
            # 各スライドの内容を抽出（最初の5スライドまで）
            max_slides = min(len(prs.slides), 5)
            
            for i, slide in enumerate(prs.slides[:max_slides]):
                ppt_text.append(f"\n[スライド {i+1}]")
                
                # スライドのタイトル
                if slide.shapes.title:
                    title_text = slide.shapes.title.text
                    ppt_text.append(f"タイトル: {title_text}")
                
                # スライド内のテキストを抽出
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape != slide.shapes.title:  # タイトル以外のテキスト
                            texts.append(shape.text)
                
                if texts:
                    ppt_text.append("内容:")
                    for text in texts:
                        ppt_text.append(f"  {text}")
                else:
                    ppt_text.append("テキスト内容なし")
            
            if len(prs.slides) > max_slides:
                ppt_text.append(f"\n...(残り {len(prs.slides) - max_slides} スライドは省略されました)")
            
            return "\n".join(ppt_text)
            
        except Exception as e:
            logger.error(f"PowerPoint抽出エラー: {str(e)}")
            return f"PowerPointファイル '{os.path.basename(file_path)}' の読み込みに失敗しました: {str(e)}"

    def get_relevant_content(self, query, max_files=None, max_chars=8000):
        """
        クエリに関連する内容を取得

        Args:
            query: 検索クエリ
            max_files: 取得する最大ファイル数
            max_chars: 取得する最大文字数

        Returns:
            関連コンテンツ（文字列）
        """
        # 最大ファイル数の設定
        if max_files is None:
            max_files = self.max_results

        # 日付を抽出（YYYY年MM月DD日）
        date_match = re.search(r'(\d{4})年(\d{1,2})月(\d{1,2})日', query)
        date_str = None

        if date_match:
            year = date_match.group(1)
            month = date_match.group(2).zfill(2)
            day = date_match.group(3).zfill(2)
            date_str = f"{year}年{month}月{day}日"
            date_pattern = f"{year}{month}{day}"
            logger.info(f"日付指定を検出: {date_str} (パターン: {date_pattern})")

        # 検索クエリからストップワードを除去
        stop_words = ["について", "とは", "の", "を", "に", "は", "で", "が", "と", "から", "へ", "より", 
                     "内容", "知りたい", "あったのか", "何", "教えて", "どのような", "どんな", "ありました",
                     "提示して", "内容は", "報告は", "報告書", "教えてください",
                     "the", "a", "an", "and", "or", "of", "to", "in", "for", "on", "by"]

        # クエリから重要な単語を抽出
        keywords = []

        # 先に日付を追加（もし存在すれば）
        if date_str:
            keywords.append(date_str)

        # その他のキーワードを追加
        for word in query.split():
            clean_word = word.strip(',.;:!?()[]{}"\'')
            if clean_word and len(clean_word) > 1 and clean_word.lower() not in stop_words:
                # 日付文字列の一部でなければ追加
                if date_str and date_str not in clean_word:
                    keywords.append(clean_word)
                elif not date_str:
                    keywords.append(clean_word)
                else:
                    pass

        # キーワードが少なすぎる場合のバックアップとして日報関連の単語を追加
        if len(keywords) < 2:
            if "日報" not in query.lower() and not any(k for k in keywords if "日報" in k):
                keywords.append("日報")

        if not keywords:
            return "検索キーワードが見つかりませんでした。具体的な日付や単語で検索してください。"

        logger.info(f"抽出されたキーワード: {keywords}")

        # ファイル検索
        search_results = self.search_files(keywords, max_results=max_files)

        if not search_results:
            keywords_str = ", ".join(keywords)
            # 日付指定がある場合は特別なメッセージ
            if date_str:
                return f"{date_str}の日報は見つかりませんでした。日付の表記が正しいか確認してください。"
            else:
                return f"キーワード '{keywords_str}' に関連するファイルは見つかりませんでした。"

        # 関連コンテンツの取得
        relevant_content = f"--- {len(search_results)}件の関連ファイルが見つかりました ---\n\n"
        total_chars = len(relevant_content)

        for i, result in enumerate(search_results):
            file_path = result.get('path')
            file_name = result.get('name')
            modified = result.get('modified', '不明')

            # ファイルの内容を読み込み
            content = self.read_file_content(file_path)

            # 各ファイルのコンテンツ用に十分なスペースを確保（均等配分）
            allowed_chars_per_file = (max_chars - total_chars) // (len(search_results) - i)
            
            # ヘッダー部分
            file_header = f"=== ファイル {i+1}: {file_name} ===\n"
            file_header += f"更新日時: {modified}\n"
            
            # プレビュー長を計算（ヘッダーの長さを考慮）
            preview_length = allowed_chars_per_file - len(file_header) - 10  # 余裕を持たせる
            if preview_length < 100:  # 最低表示文字数
                preview_length = 100
                
            # コンテンツの要約
            if len(content) > preview_length:
                preview = content[:preview_length] + "...\n"
            else:
                preview = content + "\n"
                
            file_content = file_header + preview + "\n"

            # 最大文字数をチェック
            if total_chars + len(file_content) > max_chars:
                # 制限に達した場合は切り詰め
                remaining = max_chars - total_chars - 100  # 終了メッセージ用に余裕を持たせる
                if remaining > 200:  # 最低表示文字数
                    truncated_content = file_content[:remaining] + "...\n"
                    relevant_content += truncated_content
                    total_chars += len(truncated_content)
                
                relevant_content += f"\n（残り{len(search_results) - i}件のファイルは文字数制限のため表示されません）"
                break

            relevant_content += file_content
            total_chars += len(file_content)

        return relevant_content

# 使用例
if __name__ == "__main__":
    # ロギングの設定
    logging.basicConfig(
        level=logging.DEBUG,
        format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
    )

    # インスタンス作成
    onedrive_search = OneDriveSearch()

    # 検索クエリ
    test_query = "2025年3月15日の日報内容"

    # 関連コンテンツを取得
    content = onedrive_search.get_relevant_content(test_query)

    print(f"検索結果: {content[:500]}...")  # 最初の500文字のみ表示
